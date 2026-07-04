import collections
from pptx import Presentation

def analyze():
    filename = 'Minimal_Icons 3.pptx'
    prs = Presentation(filename)
    
    num_slides = len(prs.slides)
    print(f"1. Number of slides: {num_slides}")
    
    # Analyze Slide 2 (index 1) which has the first table
    slide = prs.slides[1]
    
    table_shape = None
    floating_graphics = []
    
    for shape in slide.shapes:
        if shape.has_table:
            table_shape = shape
        else:
            if shape.name.startswith("Graphic") or shape.name.startswith("Picture") or shape.name != "Table":
                # Only count images/svgs
                if shape.shape_type == 13: # PICTURE
                    floating_graphics.append(shape.name)
                
    if table_shape is not None:
        table = table_shape.table
        rows = len(table.rows)
        cols = len(table.columns)
        print(f"2. Table dimensions per slide (based on slide 2): {rows} rows, {cols} columns")
        
        print(f"3. Floating SVG/graphic shapes on the first grid slide (Slide 2): {len(floating_graphics)}")
        
        print("4. First 10 labels from the table's first column (Slide 2):")
        count = 0
        for i, row in enumerate(table.rows):
            if count >= 10:
                break
            # the text is in the first column
            cell_text = row.cells[0].text_frame.text.replace('\n', ' ').strip()
            print(f"   Row {i+1}: {cell_text}")
            count += 1
    else:
        print("No table found on Slide 2.")

if __name__ == '__main__':
    analyze()
