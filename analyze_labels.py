import collections
from pptx import Presentation

def analyze():
    filename = 'Minimal_Icons 3.pptx'
    prs = Presentation(filename)
    
    # 4. Print the first 10 labels from the table's first column
    print("Labels:")
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for i, row in enumerate(table.rows):
                    cell_text = row.cells[0].text_frame.text.replace('\n', ' ').strip()
                    if cell_text: # only print non-empty
                        print(f"{count+1}: {cell_text}")
                        count += 1
                        if count >= 10:
                            return

if __name__ == '__main__':
    analyze()
