import collections
import json
from pptx import Presentation

def extract_icons():
    filename = 'Minimal_Icons 3.pptx'
    prs = Presentation(filename)
    
    results = []
    graphic_id = 1
    
    for slide_idx, slide in enumerate(prs.slides):
        # Find the table on this slide
        table_shape = None
        for shape in slide.shapes:
            if shape.has_table:
                table_shape = shape
                break
                
        if not table_shape:
            continue
            
        table = table_shape.table
        rows_info = []
        current_top = table_shape.top
        
        for row in table.rows:
            last_idx = len(row.cells) - 1
            label = row.cells[last_idx].text_frame.text.replace('\n', ' ').strip()
            if not label:
                label = row.cells[0].text_frame.text.replace('\n', ' ').strip()
                
            center_y = current_top + row.height / 2
            rows_info.append({
                'label': label,
                'top': current_top,
                'center_y': center_y
            })
            current_top += row.height
            
        for shape in slide.shapes:
            if shape.has_table:
                continue
                
            if shape.shape_type == 13 or shape.name.startswith("Graphic") or shape.name.startswith("Picture"):
                graphic_center_y = shape.top + shape.height / 2
                
                nearest_row = min(rows_info, key=lambda r: abs(r['center_y'] - graphic_center_y))
                label = nearest_row['label']
                
                rId = None
                try:
                    # Use local-name() to avoid namespace issues
                    svg_blips = shape._element.xpath(".//*[local-name()='svgBlip']")
                    if svg_blips:
                        rId = svg_blips[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                    else:
                        blips = shape._element.xpath(".//*[local-name()='blip']")
                        if blips:
                            rId = blips[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                except Exception as e:
                    print("Error parsing XML:", e)
                    pass
                    
                media_file = "unknown"
                if rId:
                    try:
                        part = slide.part.related_part(rId)
                        media_file = part.partname.split('/')[-1]
                    except Exception as e:
                        print("Error getting related part:", e)
                        pass
                
                if media_file == "unknown":
                    continue
                    
                results.append({
                    "id": graphic_id,
                    "file": media_file,
                    "label": label
                })
                graphic_id += 1

    with open('output.json', 'w', encoding='utf-8') as f:
        json.dump(results, f, indent=2)

    print(f"Extraction complete. Found {len(results)} graphics.")
    
if __name__ == '__main__':
    extract_icons()
