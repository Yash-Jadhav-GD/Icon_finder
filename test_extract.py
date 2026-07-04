import collections
from pptx import Presentation

def test_extract():
    filename = 'Minimal_Icons 3.pptx'
    prs = Presentation(filename)
    
    slide = prs.slides[1]
    for shape in slide.shapes:
        if shape.shape_type == 13 or shape.name.startswith("Graphic"):
            print("Found shape:", shape.name)
            svg_blips = shape._element.xpath('.//asvg:svgBlip')
            blips = shape._element.xpath('.//a:blip')
            print("SVG blips:", len(svg_blips))
            print("a:blip:", len(blips))
            
            try:
                print("XML:", shape._element.xml[:300]) # Print beginning of XML to inspect
            except:
                pass
            break

if __name__ == '__main__':
    test_extract()
